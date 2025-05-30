import os
import datetime
from pptx import Presentation
from midjourney_api import MidjourneyClient  # Hypothetical Midjourney API client

# Initialize Midjourney API client (replace with your actual API key/setup)
midjourney_client = MidjourneyClient(api_key="YOUR_API_KEY")

# Define the slides data with highly detailed visual descriptions
slides = [
    # Section 1: Introduction – The Medium is the Message
    {
        "title": "Why Software Matters",
        "text": "Software is not just a tool; it’s a medium that reshapes how we think, interact, and live.",
        "visual_description": "In a dimly lit attic, a dusty 1980s vintage computer with a bulky CRT monitor, beige keyboard, and floppy disk drive sits on a wooden table. Its screen flickers as it slowly morphs into a pulsating, realistic human brain with glowing neural circuits, set against a backdrop of faint binary code raining down, symbolizing the deep fusion of technology with human cognition and its transformative power over decades.",
        "notes": (
            "Welcome the audience and introduce the theme. Explain how software, like all media, extends our capabilities but also shapes our perceptions. "
            "Quote McLuhan: 'The medium is the message.' Highlight that software is an extension of human thought and action."
        )
    },
    {
        "title": "The Four Laws of Media",
        "text": "Every medium enhances, obsoletes, retrieves, and reverses aspects of human experience.",
        "visual_description": "A circular tetrad diagram hovers in a starry night sky, glowing softly with neon blue outlines. It’s divided into four quadrants, each with a distinct icon: a shining lightbulb for enhancement, a cracked hourglass for obsolescence, a scroll unfurling for retrieval, and a mirrored arrow for reversal. Golden arrows loop between the quadrants, pulsing as they connect, while faint text labels (e.g., 'Enhances: Connection') fade in and out, illustrating the dynamic interplay of McLuhan’s laws.",
        "notes": (
            "Introduce McLuhan’s four laws with a tetrad diagram. Explain that these laws reveal media’s hidden effects. "
            "Quote McLuhan: 'Every extension is also an amputation.' Use this to frame the discussion."
        )
    },
    {
        "title": "Today’s Mission",
        "text": "Apply the four laws to social media, YouTube, AI, and software to explore their potential for empowerment or control.",
        "visual_description": "A traveler stands at a fork in a vibrant digital landscape with a pixelated sky. The left path winds through a lush, glowing meadow with collaborative figures building together under warm sunlight, labeled 'Empowerment' on a wooden signpost. The right path descends into a shadowy, jagged ravine with towering surveillance cameras and chains, marked 'Control' in stark red letters, emphasizing the critical choice between technology’s dual futures.",
        "notes": (
            "Outline the presentation’s goal. Use a split image to show the choice between empowerment and control. "
            "Quote McLuhan: 'We shape our tools, and thereafter our tools shape us.' Set the stage for critical analysis."
        )
    },

    # Section 2: Social Media – Connecting or Isolating?
    {
        "title": "Enhancement – Amplifying Connection",
        "text": "Social media makes communication instant, global, and accessible.",
        "visual_description": "A rotating 3D globe floats in a dark void, its continents etched in gold. Bright blue, pulsating data streams arc between cities—New York, Tokyo, Lagos—forming a web of light. Zoomed-in vignettes show people typing on phones and laptops, their messages transforming into glowing lines that connect across oceans, all set against a faint hum of digital chatter, showcasing social media’s global reach.",
        "notes": (
            "Discuss how social media enhances communication. Use a world map visual to show global links. "
            "Quote McLuhan: 'The new electronic interdependence recreates the world in the image of a global village.'"
        )
    },
    {
        "title": "Obsolescence – Fading Old Bonds",
        "text": "Traditional communication methods like letters and phone calls are becoming obsolete.",
        "visual_description": "In a cozy, candlelit study, an old handwritten letter on yellowed parchment sits on an oak desk. As wind rustles through an open window, the letter crumbles into shimmering pixels that swirl upward, reassembling into a sleek smartphone glowing with notifications, set against a backdrop of fading rotary phones and typewriters, symbolizing the shift from physical to digital ties.",
        "notes": (
            "Explain how social media replaces older forms of communication. Use a fading letter visual. "
            "Quote McLuhan: 'Obsolescence never meant the end of anything, it’s just the beginning of something else.'"
        )
    },
    {
        "title": "Retrieval – The Global Village",
        "text": "Social media revives the sense of a connected, tribal community.",
        "visual_description": "In a lush digital forest under a starry sky, a crackling campfire glows orange. Around it, diverse avatars—wearing traditional garments from Japan, India, Brazil—sit on logs, their faces lit by the flames as they laugh and share animated stories that float as holograms above the fire. A faint sound of drumming ties the scene to ancient tribal roots, reimagined online.",
        "notes": (
            "Describe how social media brings back the tribal campfire experience. Use a digital campfire visual. "
            "Emphasize the return of communal interaction in a digital form."
        )
    },
    {
        "title": "Reversal – Isolation in a Connected World",
        "text": "Overuse of social media can lead to isolation and echo chambers.",
        "visual_description": "In a neon-lit city square at dusk, a lone figure in a gray hoodie stands still, head bowed. Around them, a crowd moves in a blur, each person fixated on glowing smartphone screens that cast cold light on their faces. The central figure’s shadow stretches long and empty, while faint echoes of notifications ping in the air, underscoring the paradox of isolation amid hyper-connectivity.",
        "notes": (
            "Discuss the unintended consequence of isolation. Use a lone figure visual. "
            "Quote McLuhan: 'The more you create village conditions, the more discontinuity and division and diversity.'"
        )
    },
    {
        "title": "Social Media Manifesto",
        "text": "Design platforms for authentic interaction, prioritize user control, and promote digital literacy.",
        "visual_description": "A sleek tablet floats in a minimalist white room, displaying a futuristic social media app. Its transparent interface reveals layered panels: a green 'Privacy Settings' slider, a blue 'Data Control' toggle, and a community charter written in gold text. Users’ hands adjust the settings, while a soft glow radiates from the screen, symbolizing empowerment and clarity in a user-centric design.",
        "notes": (
            "Present a manifesto for social media. Use a mock-up visual of a transparent interface. "
            "Quote McLuhan: 'There is absolutely no inevitability as long as there is a willingness to contemplate what is happening.'"
        )
    },

    # Section 3: YouTube – Democratizing or Overwhelming?
    {
        "title": "Enhancement – Access for All",
        "text": "YouTube empowers anyone to create and share video content.",
        "visual_description": "In a cluttered bedroom with posters on the walls, a teenager in a bright red hoodie films a dance video with a smartphone on a tripod. A holographic globe spins beside them, projecting a diverse crowd—kids in Paris, elders in Mumbai—cheering and commenting in real-time as colorful likes and hearts burst from the screen, celebrating the democratization of creativity.",
        "notes": (
            "Explain how YouTube enhances access to content creation. Use a smartphone filming visual. "
            "Quote McLuhan: 'The user is the content.' Highlight the democratization of media."
        )
    },
    {
        "title": "Obsolescence – End of Traditional TV",
        "text": "Scheduled broadcasts and theaters are losing relevance to on-demand streaming.",
        "visual_description": "In a dusty living room with peeling wallpaper, an old CRT TV flickers with static, a spiderweb stretching across its cracked screen. Beside it, a sleek monitor bursts to life with a vibrant YouTube homepage—thumbnails of cooking, gaming, and travel vlogs pulse with color, while the TV fades into grayscale, symbolizing the shift from rigid broadcasts to dynamic streaming.",
        "notes": (
            "Discuss how YouTube makes traditional TV obsolete. Use a dusty TV visual. "
            "Mention the shift from passive viewing to active consumption."
        )
    },
    {
        "title": "Retrieval – Oral Storytelling",
        "text": "YouTube revives the art of sharing stories and knowledge orally.",
        "visual_description": "In a virtual reality village square with cobblestone paths, a vlogger in a flowing bard’s cloak of green and gold strums a lute, narrating a tale to a circle of glowing avatars. Their faces reflect the firelight of a central torch as animated visuals of dragons and heroes rise from the story, blending medieval oral tradition with a high-tech digital stage.",
        "notes": (
            "Describe how YouTube brings back oral traditions. Use a vlogger-as-bard visual. "
            "Quote McLuhan: 'The content of any medium is always another medium.'"
        )
    },
    {
        "title": "Reversal – Information Overload",
        "text": "Too much content leads to noise, distraction, and algorithm-driven traps.",
        "visual_description": "A frantic figure in a life preserver flails in a stormy sea of swirling YouTube thumbnails—vlogs, tutorials, cat videos—each glowing and overlapping, pulling them under. Lightning cracks above as algorithm gears churn in the sky, spitting out more thumbnails, capturing the overwhelming chaos of endless content competing for attention.",
        "notes": (
            "Discuss the reversal into overload. Use a sea of thumbnails visual. "
            "Quote McLuhan: 'Too much of anything, however sweet, will always bring on the opposite.'"
        )
    },
    {
        "title": "YouTube Manifesto",
        "text": "Curate content for quality, empower users to customize algorithms, and prioritize meaningful discovery.",
        "visual_description": "A smartphone rests on a wooden table in a cozy café, its screen displaying a redesigned YouTube app. A sleek interface offers glowing filter buttons—'Educational' in blue, 'Entertainment' in red, 'Personalized' in green—while a user’s finger taps to refine their feed. Soft ambient light frames the scene, suggesting a tailored, intentional viewing experience.",
        "notes": (
            "Present a manifesto for YouTube. Use a redesigned interface visual. "
            "Encourage the audience to think about balancing freedom and focus."
        )
    },

    # Section 4: AI and Automation – Efficiency or Dependency?
    {
        "title": "Enhancement – Supercharging Productivity",
        "text": "AI automates tasks, freeing humans for creative work.",
        "visual_description": "In a gleaming high-tech lab with glass walls, a human engineer in a white coat collaborates with a silver humanoid robot. They assemble a glowing, intricate machine, passing holographic blueprints between them under bright LED lights. Tools float in midair as the robot’s precise movements complement the human’s creative input, showcasing seamless teamwork.",
        "notes": (
            "Explain how AI enhances productivity. Use a human-robot collaboration visual. "
            "Adapt McLuhan: 'The robot is an extension of our nervous system.'"
        )
    },
    {
        "title": "Obsolescence – Replacing Routine Labor",
        "text": "Manual and repetitive jobs are becoming obsolete.",
        "visual_description": "Inside a noisy factory with rusted walls, a clanking conveyor belt carries widgets past tired workers. Suddenly, the belt dissolves into streams of blue code, reforming into a pristine digital interface where AI bots in sleek pods assemble virtual goods, while the physical machines grind to a halt, fading into shadow, marking the end of manual toil.",
        "notes": (
            "Discuss how AI replaces routine labor. Use a conveyor-to-digital visual. "
            "Quote McLuhan: 'The future of work is learning.'"
        )
    },
    {
        "title": "Retrieval – The Leisure Society",
        "text": "AI revives the dream of a society focused on creativity and leisure.",
        "visual_description": "In a sun-drenched studio with open windows, an artist in a paint-splattered apron brushes vibrant colors onto a canvas, depicting a serene lake. Behind them, a team of small, whirring robots vacuums the floor, waters plants, and stirs a pot on the stove, their metallic hum blending with birdsong, illustrating a harmonious balance of automation and human creativity.",
        "notes": (
            "Describe how AI brings back the leisure society idea. Use a painting visual. "
            "Quote McLuhan: 'The future of work consists of learning a living.'"
        )
    },
    {
        "title": "Reversal – Dependency and Deskilling",
        "text": "Over-reliance on AI risks loss of skills and societal vulnerability.",
        "visual_description": "In a stark, gray room, a trembling human hand reaches for a wrench on a table but falters, weak and unsteady. Beside it, a gleaming robotic arm with red joints effortlessly twists the tool, tightening a bolt on a machine. The human hand shrinks back, its shadow fading, while the robot’s whirring grows louder, warning of skill erosion through over-dependence.",
        "notes": (
            "Discuss the reversal into dependency. Use a weakened hand visual. "
            "Adapt McLuhan: 'The price of eternal vigilance is freedom.'"
        )
    },
    {
        "title": "AI Manifesto",
        "text": "Build transparent, human-augmenting AI with oversight and ethical controls.",
        "visual_description": "A futuristic dashboard in a control room glows with a transparent AI interface. Sliders labeled 'Transparency' (green), 'Fairness' (blue), and 'Accountability' (yellow) adjust with a human hand, while a 'Human Oversight' button pulses red. Holographic figures of diverse people nod in approval, bathed in soft light, promoting an ethical partnership with AI.",
        "notes": (
            "Present a manifesto for AI. Use an ethical AI interface visual. "
            "Encourage the audience to consider AI as an enhancer, not a replacement."
        )
    },

    # Section 5: Software Generally – Empowerment or Complexity?
    {
        "title": "Enhancement – Extending Human Capabilities",
        "text": "Software amplifies creativity, productivity, and problem-solving.",
        "visual_description": "A coder in a sleek black superhero suit stands atop a skyscraper at night, city lights twinkling below. Their utility belt glows with tools—code editors (blue), debuggers (red), version control (green)—as they type midair, conjuring a holographic bridge of binary code that spans two buildings, symbolizing software’s power to solve real-world challenges.",
        "notes": (
            "Explain how software enhances human capabilities. Use a toolbelt visual. "
            "Quote McLuhan: 'We become what we behold.'"
        )
    },
    {
        "title": "Obsolescence – Outdating Analog Tools",
        "text": "Typewriters, ledgers, and physical media are being replaced by digital solutions.",
        "visual_description": "In a cluttered office with wooden shelves, an antique typewriter clacks its last letter before trembling and crumbling into a cascade of black-and-white binary digits. The digits swirl upward, reforming into a slim silver laptop on the desk, its screen glowing with a spreadsheet, as old ledgers and film reels dissolve into dust behind it, marking the digital takeover.",
        "notes": (
            "Discuss how software makes analog tools obsolete. Use a typewriter-to-code visual. "
            "Highlight the transition to digital workflows."
        )
    },
    {
        "title": "Retrieval – Craftsmanship",
        "text": "Software development mirrors the artisanal creation of intricate systems.",
        "visual_description": "In a smoky medieval forge, a coder in a leather apron hammers a glowing keyboard on an anvil, sparks flying with each strike. Lines of neon-green code burst forth, shaping into a digital sword that pulses with energy, while a roaring fire casts shadows of ancient tools, blending the meticulous craft of blacksmithing with modern programming artistry.",
        "notes": (
            "Describe how software retrieves craftsmanship. Use a coder-as-blacksmith visual. "
            "Emphasize the skill and care in coding as a modern craft."
        )
    },
    {
        "title": "Reversal – Complexity and Fragility",
        "text": "Overcomplicated software leads to bugs, breaches, and loss of control.",
        "visual_description": "A harried developer in a dimly lit basement sits at a desk, engulfed by a towering, writhing 3D web of red code strings, blinking error alerts, and spinning debugging icons. Sweat beads on their forehead as the web tightens, a crashed program flashing on their screen, depicting the suffocating chaos of software gone awry.",
        "notes": (
            "Discuss the reversal into complexity. Use a tangled code visual. "
            "Quote McLuhan: 'The medium is the massage.' Note how complexity can overwhelm users."
        )
    },
    {
        "title": "Software Manifesto",
        "text": "Prioritize simplicity, security, and user empowerment in software design.",
        "visual_description": "In a bright, airy studio, a computer screen displays a minimalist software app with a clean white background. Icons for 'Open Source' (a green tree), 'Collaboration' (blue hands), and 'Transparency' (a yellow eye) orbit a central interface, where a user drags simple widgets into place, crafting a tool under warm sunlight, embodying intuitive and empowering design.",
        "notes": (
            "Present a manifesto for software. Use a clean interface visual. "
            "Encourage the audience to keep software as a tool for empowerment."
        )
    },

    # Section 6: Manifestos for a Human-Centered Future
    {
        "title": "The Big Picture",
        "text": "Software can empower or control—our choices define its path.",
        "visual_description": "At a crossroads in a shimmering digital desert, a lone figure in white gazes at two paths under a twilight sky. The left leads to a radiant field where diverse people in colorful robes build glowing structures together, laughing. The right descends into a dark city of steel towers, red-eyed cameras scanning silent citizens, contrasting freedom with oppression in software’s future.",
        "notes": (
            "Summarize the presentation’s theme. Use a fork in the road visual. "
            "Quote McLuhan: 'There is absolutely no inevitability as long as there is a willingness to contemplate what is happening.'"
        )
    },
    {
        "title": "Unified Manifesto",
        "text": "Design software for human agency, transparency, and ethics, guided by the four laws.",
        "visual_description": "In a grand architect’s studio, a blueprint unfurls on a wooden table, bathed in golden lamplight. It depicts a majestic building with four cornerstones, each engraved with McLuhan’s laws—enhancement (a torch), obsolescence (a broken chain), retrieval (a quill), reversal (a mirror)—while workers in hard hats lay bricks of code, symbolizing a foundation for ethical software.",
        "notes": (
            "Present a unified manifesto. Use a blueprint visual. "
            "Stress the importance of intentional design for a humane digital future."
        )
    },
    {
        "title": "Call to Action",
        "text": "Collaborate to shape software’s future—propose one empowering feature or policy.",
        "visual_description": "On a windswept cliff overlooking a digital sea, a diverse team—scientists, artists, coders in vibrant attire—builds a bridge of shimmering binary code and data streams. They hammer and weld, connecting a fractured landmass under a hopeful orange sky, their tools sparking as they unite to span the digital divide with purpose.",
        "notes": (
            "Challenge the audience to take action. Use a bridge-building visual. "
            "Quote McLuhan: 'The medium is the message.' Ask: What message will your software send?"
        )
    },

    # Section 7: Q&A and Closing Reflections
    {
        "title": "Open Discussion",
        "text": "How will you rethink software’s role in your work or life?",
        "visual_description": "In a cosmic void, a massive, glowing question mark pulses in gold at the center. Orbiting it are devices—smartphones, laptops, smartwatches—linked by thin, silver threads of light, spinning slowly. Faint voices murmur questions as the devices hum, inviting reflection on technology’s personal impact in an interconnected world.",
        "notes": (
            "Invite questions and reflections. Use a question mark visual. "
            "Encourage the audience to consider their relationship with software."
        )
    },
    {
        "title": "Closing – The Medium is Ours to Shape",
        "text": "Software’s potential is vast, but ours to mold. Let’s choose empowerment.",
        "visual_description": "A golden sunrise breaks over a futuristic city nestled in a green valley, where sleek towers blend with trees. People in flowing robes use tablets and holograms thoughtfully, tending gardens as drones hum softly overhead. The sky glows with possibility, reflecting a balanced, hopeful future shaped by conscious tech choices.",
        "notes": (
            "Conclude with a hopeful message. Use a sunrise visual. "
            "Quote McLuhan: 'We look at the present through a rearview mirror. We march backwards into the future.' Urge the audience to look forward."
        )
    }
]

def generate_gif(prompt, slide_num):
    """
    Generate an animated GIF using the Midjourney API based on the provided prompt.
    Ensure the GIF is in portrait orientation.
    
    Args:
        prompt (str): The visual description to use as the API prompt.
        slide_num (int): Slide number for naming the output file.
    
    Returns:
        str: Path to the generated GIF, or None if generation fails.
    """
    try:
        # Create a directory for temporary GIFs if it doesn’t exist
        os.makedirs("generated_gifs", exist_ok=True)
        
        # Set parameters for portrait orientation and animation
        response = midjourney_client.generate_image(
            prompt=prompt,
            aspect_ratio="9:16",  # Portrait orientation (adjust as needed)
            animation=True,       # Request an animated GIF
            output_format="gif"
        )
        
        # Assume the API returns a file path or URL; adjust based on actual API response
        gif_path = f"generated_gifs/slide_{slide_num}.gif"
        # Simulate saving the GIF locally (replace with actual download logic if needed)
        with open(gif_path, "wb") as f:
            f.write(response['gif_data'])  # Hypothetical response field
        
        return gif_path
    except Exception as e:
        print(f"Error generating GIF for slide {slide_num}: {e}")
        return None

def create_slide(prs, slide_info, i):
    """
    Create a slide based on provided slide information.
    
    Args:
        prs (Presentation): PowerPoint presentation object.
        slide_info (dict): Dictionary containing slide details.
        i (int): Slide index for naming purposes.
    """
    title = slide_info["title"]
    text = slide_info["text"]
    visual_description = slide_info["visual_description"]
    notes = slide_info["notes"]

    # Use 'Two Content' layout (index 3) for slides with images
    slide_layout = prs.slide_layouts[3]  # Assumes layout 3 is 'Two Content'
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title

    # Add text to the left content placeholder
    left_placeholder = slide.placeholders[1]
    left_placeholder.text = text

    # Generate and add GIF to the right content placeholder
    gif_path = generate_gif(visual_description, i)
    if gif_path and os.path.exists(gif_path):
        right_placeholder = slide.placeholders[2]
        left = right_placeholder.left
        top = right_placeholder.top
        width = right_placeholder.width
        height = right_placeholder.height
        slide.shapes.add_picture(gif_path, left, top, width, height)

    # Add presenter's notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = notes

def create_presentation(slides, base_filename="presentation"):
    """
    Generate a PowerPoint presentation based on provided slide information.
    
    Args:
        slides (list): List of dictionaries containing slide details (title, text, visual_description, notes).
        base_filename (str): Base name for the output file (default: 'presentation').
    """
    # Initialize a new presentation
    prs = Presentation()

    for i, slide_info in enumerate(slides):
        create_slide(prs, slide_info, i)

    # Generate filename with timestamp
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"{base_filename}_{timestamp}.pptx"
    prs.save(filename)
    print(f"Presentation saved as {filename}")

    # Clean up generated GIFs
    for i in range(len(slides)):
        gif_path = f"generated_gifs/slide_{i}.gif"
        if os.path.exists(gif_path):
            os.remove(gif_path)
    if os.path.exists("generated_gifs") and not os.listdir("generated_gifs"):
        os.rmdir("generated_gifs")

if __name__ == "__main__":
    create_presentation(slides)