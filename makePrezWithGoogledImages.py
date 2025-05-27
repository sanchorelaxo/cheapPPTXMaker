import os
import datetime
import json
from pptx import Presentation
from icrawler.builtin import GoogleImageCrawler
from PIL import Image

# Load slides from external JSON file
with open("slides.json", "r", encoding="utf-8") as f:
    slides = json.load(f)
  

def fetch_image(keywords, slide_num):
    """
    Fetch an image using GoogleImageCrawler based on the provided keywords.
    Prioritize portrait-oriented images (height > width). If no portrait image is found, use the first image.
    """
    # Set up image crawler to fetch large images
    crawler = GoogleImageCrawler(storage={'root_dir': f'images/slide_{slide_num}'})
    try:
        crawler.crawl(keyword=keywords, max_num=5, filters={'size': 'large'})
    except Exception as e:
        print(f"Image crawl failed with full query: {e}\nFalling back to shorter query.")
        short_query = " ".join(keywords.split()[:5])
        try:
            crawler.crawl(keyword=short_query, max_num=5, filters={'size': 'large'})
        except Exception as e2:
            print(f"Image crawl failed with short query as well: {e2}")

    # Find the first portrait-oriented image (height > width)
    image_path = None
    downloaded_files = os.listdir(f'images/slide_{slide_num}') if os.path.exists(f'images/slide_{slide_num}') else []
    for file in downloaded_files:
        try:
            img = Image.open(os.path.join(f'images/slide_{slide_num}', file))
            width, height = img.size
            if height > width:
                image_path = os.path.join(f'images/slide_{slide_num}', file)
                break
        except Exception:
            continue  # Skip files that can't be opened

    # If no portrait image is found, use the first downloaded image as fallback
    if image_path is None and downloaded_files:
        image_path = os.path.join(f'images/slide_{slide_num}', downloaded_files[0])

    return image_path


def create_slide(prs, slide_info, i):
    """
    Create a slide based on provided slide information.
    """
    title = slide_info["title"]
    text = slide_info["text"]
    keywords = slide_info["keywords"]
    notes = slide_info["notes"]

    # Use 'Two Content' layout (index 3) for slides with images
    slide_layout = prs.slide_layouts[3]  # Assumes layout 3 is 'Two Content'
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title

    # Add text to the left content placeholder
    left_placeholder = slide.placeholders[1]
    left_placeholder.text = text

    # Fetch and add image to the right content placeholder
    image_path = fetch_image(keywords, i)
    if image_path:
        right_placeholder = slide.placeholders[2]
        left = right_placeholder.left
        top = right_placeholder.top
        width = right_placeholder.width
        height = right_placeholder.height
        # Always convert to PNG for pptx compatibility
        from uuid import uuid4
        img = Image.open(image_path)
        converted_path = f"{image_path}.{uuid4().hex}.converted.png"
        img.save(converted_path, 'PNG')
        slide.shapes.add_picture(converted_path, left, top, width, height)
        if os.path.exists(converted_path):
            os.remove(converted_path)

    # Add presenter's notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = notes


def create_presentation(slides, base_filename="presentation"):
    """
    Generate a PowerPoint presentation based on provided slide information.
    
    Args:
        slides (list): List of dictionaries containing slide details (title, text, keywords, notes).
        base_filename (str): Base name for the output file (default: 'presentation').
    """
    # Initialize a new presentation
    prs = Presentation()

    for i, slide_info in enumerate(slides):
        create_slide(prs, slide_info, i)

    # Generate filename with timestamp
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"{base_filename}_{timestamp}.pptx"
    prs.save(filename)
    print(f"Presentation saved as {filename}")

    # Clean up downloaded images
    for i in range(len(slides)):
        slide_dir = f'images/slide_{i}'
        if os.path.exists(slide_dir):
            for file in os.listdir(slide_dir):
                os.remove(os.path.join(slide_dir, file))
            os.rmdir(slide_dir)


if __name__ == "__main__":
    create_presentation(slides)